from docx import Document
from pptx import Presentation
import os
from pptx.util import Pt


def lire_texte_docx(nom_fichier):
    doc = Document(nom_fichier)
    texte = []
    for paragraphe in doc.paragraphs:
        texte.append(paragraphe.text)
    return "\n".join(texte)


def premiere_ligne_docx(nom_fichier):
    doc = Document(nom_fichier)
    if doc.paragraphs:
        premiere_ligne = doc.paragraphs[0].text.split()  # Diviser la ligne en mots
        return premiere_ligne
    else:
        return ["Le", "document", "est", "vide."]


def extraire_texte_apres_ligne(nom_fichier, texte_debut):
    doc = Document(nom_fichier)
    texte_trouve = False
    texte_suivant = None

    # Recherche du texte de départ pour déterminer où commencer l'extraction
    for paragraphe in doc.paragraphs:
        if texte_trouve:
            texte_suivant = paragraphe.text
            break  # Arrêter l'extraction après avoir trouvé la première ligne suivant le texte spécifié
        if texte_debut in paragraphe.text:
            texte_trouve = True

    return texte_suivant


def remplacer_texte_premier_slide(presentation, nouveaux_mots):
    premier_slide = presentation.slides[0]

    mots_a_remplacer = ["ANDRO", "00", "VOLANA", "taona"]

    for forme in premier_slide.shapes:
        if not forme.has_text_frame:
            continue
        for paragraphe in forme.text_frame.paragraphs:
            for run in paragraphe.runs:
                for i, mot in enumerate(mots_a_remplacer):
                    if mot in run.text:
                        run.text = run.text.replace(mot, nouveaux_mots[i])


def remplacer_texte_deuxieme_slide(presentation, texte_a_remplacer, nouveau_texte):
    deuxieme_slide = presentation.slides[1]  # Deuxième slide (index 1)

    for forme in deuxieme_slide.shapes:
        if not forme.has_text_frame:
            continue
        for paragraphe in forme.text_frame.paragraphs:
            for run in paragraphe.runs:
                if texte_a_remplacer in run.text:
                    run.text = run.text.replace(texte_a_remplacer, nouveau_texte)


def remplacer_texte_powerpoint(
    nom_fichier_powerpoint, nouveaux_mots, texte_debut, nouveau_texte
):
    presentation = Presentation(nom_fichier_powerpoint)

    # Remplacer sur le premier slide
    remplacer_texte_premier_slide(presentation, nouveaux_mots)

    # Remplacer sur le deuxième slide
    remplacer_texte_deuxieme_slide(presentation, texte_debut, nouveau_texte)

    chemin_nouveau_fichier = "fafana.pptx"
    presentation.save(chemin_nouveau_fichier)
    return chemin_nouveau_fichier


def extraire_trois_lignes_apres_mot(nom_fichier, texte_debut, nombre_lignes=3):
    doc = Document(nom_fichier)
    lignes = []
    texte_trouve = False

    for paragraphe in doc.paragraphs:
        if texte_debut in paragraphe.text:
            texte_trouve = True
            continue
        if texte_trouve and len(lignes) < nombre_lignes:
            lignes.append(paragraphe.text)

    return lignes


def remplacer_texte_troisieme_slide(presentation, texte_a_remplacer, nouveau_texte,nouvelle_taille=115):
    troisieme_slide = presentation.slides[2]  # Troisième slide (index 2)

    for forme in troisieme_slide.shapes:
        if not forme.has_text_frame:
            continue
        for paragraphe in forme.text_frame.paragraphs:
            for run in paragraphe.runs:
                if texte_a_remplacer in run.text:
                    run.text = run.text.replace(texte_a_remplacer, (nouveau_texte))
                    run.font.size = Pt(nouvelle_taille)  # Réglez la taille de la police

    return presentation  # Retourne la présentation modifiée


def creer_fichier(chemin_fichier_docx, chemin_fichier_powerpoint):
    texte_docx = lire_texte_docx(chemin_fichier_docx)
    premiere_ligne = premiere_ligne_docx(chemin_fichier_docx)
    print(texte_docx)
    print(premiere_ligne)

    texte_a_remplacer_deuxieme_slide = "SALAMO VERSET"

    nouveau_fichier = remplacer_texte_powerpoint(
        chemin_fichier_powerpoint,
        premiere_ligne[:4],
        texte_a_remplacer_deuxieme_slide,
        extraire_texte_apres_ligne(chemin_fichier_docx, "FIARAHABANA APOSTOLIKA"),
    )

    texte_apres_soratra_masina = extraire_trois_lignes_apres_mot(
        chemin_fichier_docx, "Soratra Masina"
    )
    resultat_transforme = []

    for ligne in texte_apres_soratra_masina:
        mots = ligne.split()  # Diviser la chaîne en mots
        print(mots)
        if len(mots) >= 7:  # Si la chaîne a au moins sept mots
            nouveau_mot = f"{mots[0]} {mots[1][:3]}. {mots[2]} : {mots[4]} - {mots[6]}"  # Créer la nouvelle chaîne
        else:
            nouveau_mot = f"{mots[0][:3]}. {mots[1]} : {mots[3]} - {mots[5]}"  # Sinon, créer la nouvelle chaîne
        resultat_transforme.append(nouveau_mot)

    presentation = Presentation(nouveau_fichier)
    # Création du texte avec saut de ligne après chaque valeur
    texte_avec_saute_ligne = "\n".join(resultat_transforme)
    print("texte_avec_saute_ligne", texte_avec_saute_ligne)

    remplacer_texte_troisieme_slide(presentation, "VERSET", texte_avec_saute_ligne)
    nouveau_fichier_final = "nouveau_fichier_final.pptx"
    presentation.save(nouveau_fichier_final)
    os.startfile(nouveau_fichier_final)


# Remplacez 'chemin/vers/votre/fichier.docx' par le chemin de votre fichier .docx
chemin_fichier = 'fandaharana/ALAHADY 05 NOVAMBRA.docx'
# chemin_fichier = "fandaharana/ALAHADY 29 OKTOBRA.docx"

chemin_fichier_powerpoint = "vert.pptx"

creer_fichier(chemin_fichier, chemin_fichier_powerpoint)
