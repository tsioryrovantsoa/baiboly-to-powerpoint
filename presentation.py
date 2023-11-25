import os
from pptx import Presentation
from pptx.util import Pt
from function import *
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.util import Inches
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from dotenv import load_dotenv

load_dotenv()

def creer(idBoky, toko, andininydeb, andininyfin):
    pres = Presentation(os.getenv("TEMPLATE_PATH"))
    slide0 = pres.slides[0]
    title_shapes = [shape for shape in slide0.shapes if shape.has_text_frame]
    title = [
        shape for shape in title_shapes if shape.has_text_frame and shape.text == "boky"
    ]

    conn = create_connection()
    with conn:
        title[0].text = getBoky(conn, idBoky, toko, andininydeb, andininyfin)
        format_title(slide0)
        rows = getVerser(conn, idBoky, toko, andininydeb, andininyfin)

    for row in rows:
        create_verset_slides(row, pres)

    fichier = os.getenv("TEST_PATH")
    pres.save(fichier)
    return fichier


def format_title(slide0):
    title_para = slide0.shapes.title.text_frame.paragraphs[0]
    title_para2 = slide0.shapes.title.text_frame.paragraphs[1]
    title_para.font.name = "Alégre Sans NC"
    size = (
        Pt(100)
        if title_para.text == "Tonon-kiran'i Solomona "
        else Pt(134)
        if title_para.text == "Asan'ny Apostoly "
        else Pt(170)
        if title_para.text == "Deotoronomia "
        else Pt(175)
        if title_para.text == "I Tesaloniana "
        else Pt(175)
        if title_para.text == "II Tesaloniana "
        else Pt(185)
    )
    title_para.font.size = size
    title_para2.font.name = "Alégre Sans NC"
    title_para2.font.size = Pt(185)
    title_para.line_spacing = Pt(12)
    title_para.font.color.rgb = RGBColor(255, 255, 0)
    title_para2.font.color.rgb = RGBColor(255, 255, 0)


def create_verset_slides(row, pres):
    # row[0] : toko, row[1] : verset
    verset = row[1]
    # modele 5 du nouvelle diapositive
    first_slide_layout = pres.slide_layouts[5]

    # compter le nombre de mot verset.split()
    # si nombre de mot superieur a 13
    if len(verset.split()) > 13:
        # ponctuations pour separer des slides par ordre de priorite
        ponctuations = [",", ";", "!", "?"]
        # division = 0
        zaraina = 0

        for ponctuation in ponctuations:
            # si ponction present dans le verset
            if ponctuation in verset:
                # compter le nombre de ponctuation dans le verset et stoper le boucle
                zaraina = verset.count(ponctuation)
                break

        # verset diviser par le ponctualion par le division
        # versetoff = verset.split(ponctuation, zaraina)
        versetoff = split_verse(verset, ponctuations)
        # numero de verset
        base_text = str(row[0])

        for i, x in enumerate(versetoff):
            # i l'indice du versetoff
            # x le versetoff
            # si verset non vide
            if x:
                # creation de nouvelle diapositive
                slide = pres.slides.add_slide(first_slide_layout)
                # text = base_text + verset si debut de diapo c'est a dire indice 0 si non text = verset
                text = f"{base_text} {x}" if i == 0 else x
                # propriete du slide : slide, text, le taille de police
                set_slide_properties(slide, text, size=Pont_size(x))
    # si nombre inferieur a 13
    else:
        # creation de nouvelle diapositive
        slide = pres.slides.add_slide(first_slide_layout)
        # text = base_text + verset
        text = f"{row[0]} {row[1]}"
        # propriete du slide : slide, text, le taille de police
        set_slide_properties(slide, text, size=Pont_size(row[1]))


def split_verse(verset, ponctuations):
    parts = []

    # Initialise le premier morceau avec une chaîne vide
    current_part = ""

    # Parcourt chaque caractère du verset
    for char in verset:
        # Si le caractère est une ponctuation à diviser
        if char in ponctuations:
            # Ajoute le caractère à la fin du morceau actuel
            current_part += char
            # Ajoute le morceau actuel à la liste des parties
            parts.append(current_part)
            # Réinitialise le morceau actuel à une chaîne vide
            current_part = ""
        else:
            # Ajoute le caractère au morceau actuel
            current_part += char

    # Ajoute le dernier morceau (après la dernière ponctuation)
    parts.append(current_part)

    return parts


def set_slide_properties(slide, text, size):
    # ajouter text dans le diapo
    slide.shapes.title.text = text
    slide.shapes.title.width = Inches(10)
    slide.shapes.title.height = Inches(7.5)
    slide.shapes.title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    slide.shapes.title.text_frame.paragraphs[0].font.name = "Helvetica Inserat LT Std"
    slide.shapes.title.text_frame.paragraphs[0].font.size = size


def Pont_size(text):
    # Police 72 si nombre de mot dans le texte = 12 si non 80
    return Pt(72) if len(text.split()) == 12 else Pt(80)


def openFile(pres):
    fileName = pres
    os.system("start " + fileName)
